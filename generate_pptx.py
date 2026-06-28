"""
Generate FPL AI thesis presentation as .pptx
Dark theme: bg #080c14, title green #00e676, body white, accent purple #7c3aed
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
BG       = RGBColor(0x08, 0x0c, 0x14)
GREEN    = RGBColor(0x00, 0xe6, 0x76)
WHITE    = RGBColor(0xdd, 0xee, 0xff)
MUTED    = RGBColor(0x6a, 0x88, 0xaa)
PURPLE   = RGBColor(0x7c, 0x3a, 0xed)
AMBER    = RGBColor(0xff, 0xb3, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

slides_data = [
    {
        "title": "The Premier League & Fantasy Football",
        "bullets": [
            "England's top football division — 20 clubs, 38 gameweeks per season",
            "FPL: pick 15 real players, £100m budget, score points from real match stats",
            "Goals, assists, clean sheets, minutes played all convert to points",
            "11 million+ managers compete worldwide every season",
        ],
        "notes": (
            "Before we get into the AI, let me quickly set the scene. The Premier League is "
            "England's top football league — 20 teams, 38 weeks of football. Fantasy Premier League "
            "is the official game where you manage a virtual squad of real players. You have a £100 "
            "million budget, you pick 15 players, and every week they score points based on what "
            "actually happens on the pitch. 11 million people play this globally."
        ),
    },
    {
        "title": "Why is FPL Hard?",
        "bullets": [
            "700+ players to pick from — budget, position & club limits make it combinatorial",
            "Injuries, rotation, form, fixture difficulty change every single week",
            "Chips add strategic depth: Wildcard, Triple Captain, Bench Boost, Free Hit",
            "Average manager scores ~50 pts/GW — consistent outperformance is rare",
        ],
        "notes": (
            "So why is this hard? On the surface it looks simple — pick good players. But there are "
            "hundreds of constraints. You have 700+ players, a tight budget, position rules, a max of "
            "3 players from any one club. The number of valid squad combinations runs into the billions. "
            "On top of that, everything changes week to week. Most managers finish the season around "
            "average. Consistent top performance requires processing more data than any human can handle."
        ),
    },
    {
        "title": "What Are We Building?",
        "bullets": [
            "An autonomous FPL manager — makes weekly decisions entirely from data",
            'Research question: "Can a machine outperform the average FPL manager using only public data?"',
            "3-layer hybrid AI:  Predict  →  Decide  →  Explain",
            "Blind test constraint: zero 2025-26 data in training — GW1 is a true test",
        ],
        "notes": (
            "The question we're trying to answer is: can a machine manage an FPL team better than a "
            "human, using only publicly available data? We built a 3-layer system. First layer predicts "
            "how many points each player will score. Second layer mathematically selects the best squad. "
            "Third layer uses an AI language model to explain every decision. GW1 was a completely blind test."
        ),
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Layer 1 — Predict:  LightGBM forecasts points for all 700+ players before each GW",
            "Layer 2 — Decide:  ILP optimizer selects optimal squad under all FPL constraints",
            "Layer 3 — Explain:  Claude LLM explains every pick in natural language",
            "Intel Pipeline: live injury data, press conferences & Gemini AI feed into Layer 2",
        ],
        "notes": (
            "Here's the full system. Data comes in from multiple sources — the FPL API, historical "
            "match data, xG statistics, press conferences. That data gets engineered into features and "
            "fed into four LightGBM models, one per position. Those predictions go into an ILP optimizer "
            "which solves squad selection as a mathematical problem. Finally Claude explains the decisions."
        ),
    },
    {
        "title": "Data Pipeline",
        "bullets": [
            "FPL API: 700+ players, prices, injuries, ownership, fixture difficulty ratings",
            "Vaastav: 6 seasons of historical GW data — 51,000+ player-gameweek rows",
            "Understat: expected goals (xG) and expected goals against (xGA) per team",
            "FBref: previous-league stats for new Premier League signings (9 leagues)",
            "Press conferences: pre-deadline injury & availability signals scraped automatically",
        ],
        "notes": (
            "The foundation is data. We pull from 5 main sources. The FPL API gives live player data. "
            "Vaastav gives 6 seasons of historical gameweek statistics back to 2019. Understat gives "
            "expected goals data which is more predictive than actual goals. FBref gives stats on new "
            "signings from foreign leagues. And we scrape press conferences before each deadline to catch "
            "injury news early."
        ),
    },
    {
        "title": "LightGBM Prediction Model",
        "bullets": [
            "4 separate models — one per position (GK, DEF, MID, FWD)",
            "Walk-forward CV: train on seasons 1-N, validate on N+1 — no data leakage",
            "Key features: rolling form, team xG/xGA, fixture difficulty, price, ownership",
            "600+ hyperparameter trials via Bayesian optimisation (Optuna TPE)",
        ],
        "notes": (
            "The prediction layer uses LightGBM — a gradient boosting model known for speed and accuracy "
            "on tabular data. We train 4 separate models, one per position, because a goalkeeper and a "
            "striker have completely different scoring patterns. We validate using walk-forward "
            "cross-validation — always training on past seasons, testing on the next one. Hyperparameters "
            "were tuned using Bayesian optimisation across 600+ trials."
        ),
    },
    {
        "title": "ILP Optimizer",
        "bullets": [
            "Integer Linear Programming — mathematically guarantees the optimal squad",
            "Hard constraints: £100m budget, position limits (2-5-5-3), max 3 per club",
            "Intel penalties: injured or rotation-risk players get prediction discounts",
            "Chip automation: TC, BB, WC, FH triggered by learned thresholds",
        ],
        "notes": (
            "Once we have predictions, we need to pick the best squad. We use Integer Linear Programming — "
            "this guarantees the globally optimal squad given the constraints. It's not heuristic, it's exact. "
            "Intel penalties discount injured or rotation-risk players before optimisation. Chips are "
            "triggered automatically based on thresholds learned during hyperparameter search."
        ),
    },
    {
        "title": "Intel Pipeline",
        "bullets": [
            "Runs automatically before each GW deadline — 7 stages",
            "Scrapes Fantasy Football Scout press conferences for injury keywords",
            "Scores every player: availability 0–100% and rotation risk 0–100%",
            "Gemini 2.5 Flash generates captain picks, differentials & risk warnings",
            "Signals become prediction multipliers fed directly into the ILP",
        ],
        "notes": (
            "The intel pipeline is what separates this from a purely historical model. Before each deadline "
            "it automatically fetches live FPL data, scrapes press conference articles to detect injuries, "
            "scores every player on availability and rotation risk, and calls Gemini AI for a pre-deadline "
            "report. All of this feeds into the optimizer as adjusted predictions."
        ),
    },
    {
        "title": "LLM Explanation Layer",
        "bullets": [
            "Claude Sonnet explains every squad selection post-simulation",
            "Per-player: model rank in position group, fixture difficulty, price value",
            "Captain analysis: why this player over the alternatives — did it pay off?",
            "Transparent AI — every decision has a reason, not a black box",
        ],
        "notes": (
            "The third layer is about transparency. After the system makes its decisions, Claude Sonnet "
            "analyses every pick and writes an explanation. For each player it explains where they ranked "
            "in the model's predictions, what their fixture looked like, and whether the pick paid off. "
            "This means the system can justify its decisions — it's not just a number that spits out a squad."
        ),
    },
    {
        "title": "Live Dashboard",
        "bullets": [
            "Flask web app — real-time view of the full 38-GW season",
            "Season overview: cumulative points chart, GW bars, chips timeline",
            "Per-GW pitch view: squad layout, actual vs predicted, transfers, bank & FT",
            "Click any player → Claude's explanation of why they were selected",
        ],
        "notes": (
            "Everything comes together in this dashboard. You can see the full season gameweek by gameweek. "
            "The season overview shows cumulative points climbing to 2474. You can navigate to any gameweek, "
            "see the exact squad the system picked, the actual scores, the transfers made, and click any "
            "player to read Claude's explanation of why they were selected."
        ),
    },
    {
        "title": "Results",
        "bullets": [
            "2,474 points across 38 gameweeks — 65.1 pts/GW average",
            "Global rank: top ~250 out of 11 million+ managers  (top 0.002%)",
            "Only 2 transfer hits all season: -4 pts GW8, -4 pts GW21",
            "Chips timed optimally: TC×2, BB×2, WC×1, FH×1 — none wasted",
        ],
        "notes": (
            "The results. Across the full 38-gameweek 2025-26 season the system scored 2,474 points — "
            "averaging 65 points per gameweek. That places it in the top ~250 managers globally out of "
            "over 11 million. The top score in the entire game was 2,582. We finished in the top 0.002%. "
            "The system took only two transfer hits the entire season, and all six chips were used."
        ),
    },
    {
        "title": "Key Findings & Limitations",
        "bullets": [
            "LightGBM outperformed XGBoost consistently across all positions",
            "Bayesian search (Optuna) found +178 pts improvement over default configuration",
            "Limitation: press scraper misses inline injury mentions (e.g. Newcastle players)",
            "Limitation: no sell-buyback penalty — ILP can re-buy a sold player next GW",
        ],
        "notes": (
            "A few key findings. LightGBM dominated XGBoost in every experiment. Bayesian hyperparameter "
            "search found configurations nearly 180 points better than the defaults — showing that tuning "
            "matters significantly. On limitations — the press conference scraper has a known blind spot "
            "for inline mentions. And the ILP has no memory between weeks so it can theoretically sell and "
            "re-buy the same player."
        ),
    },
    {
        "title": "Conclusion & Future Work",
        "bullets": [
            "A 3-layer hybrid AI reaches top 0.002% FPL performance using only public data",
            "Fully automated: data collection → prediction → optimisation → explanation",
            "Future: live season deployment, fatigue features, RL for transfer strategy",
            '"The best FPL managers combine data and intuition — this system handles the data."',
        ],
        "notes": (
            "To conclude — we built a fully automated FPL management system that finishes in the global "
            "top 250 using nothing but publicly available data and open-source machine learning. Every "
            "decision is explainable, every component is modular, and the system ran the entire 2025-26 "
            "season without human intervention. Thank you."
        ),
    },
]


def set_bg(slide, prs):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_slide(prs, title_text, bullets, notes_text):
    layout = prs.slide_layouts[6]  # blank layout
    slide  = prs.slides.add_slide(layout)
    set_bg(slide, prs)

    W = SLIDE_W
    H = SLIDE_H

    # Title bar background strip
    from pptx.util import Emu
    strip = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        W, Inches(1.3)
    )
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(0x0f, 0x16, 0x23)
    strip.line.fill.background()

    # Green accent bar (left edge)
    accent = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(0.06), Inches(1.3)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = GREEN
    accent.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.2), Inches(0.15),
        Inches(12.9), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.bold   = True
    run.font.size   = Pt(32)
    run.font.color.rgb = GREEN
    run.font.name   = "Calibri"

    # Bullet area
    body_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.5),
        Inches(12.5), Inches(5.7)
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = "▸  " + bullet
        run.font.size  = Pt(20)
        run.font.color.rgb = WHITE
        run.font.name  = "Calibri"

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text

    return slide


def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in slides_data:
        add_slide(prs, s["title"], s["bullets"], s["notes"])

    out = r"c:\Users\Andrej\Desktop\FPL_AI_Presentation.pptx"
    prs.save(out)
    print("Saved: " + out)
    print("Slides: " + str(len(slides_data)))


if __name__ == "__main__":
    main()
